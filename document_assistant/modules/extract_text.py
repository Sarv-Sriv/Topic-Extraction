import fitz
from docx import Document
from pptx import Presentation
import os


def clean_text(text):

    return " ".join(text.split())


def extract_pdf(path):

    extracted = []

    doc = fitz.open(path)

    for page_num in range(len(doc)):

        text = doc[page_num].get_text()

        text = clean_text(text)

        if len(text.strip()) < 50:
            continue

        extracted.append({
            "text": text,
            "source": os.path.basename(path),
            "page": page_num + 1
        })

    return extracted


def extract_docx(path):

    doc = Document(path)

    text = "\n".join(
        [p.text for p in doc.paragraphs]
    )

    text = clean_text(text)

    return [{
        "text": text,
        "source": os.path.basename(path),
        "page": 1
    }]


def extract_pptx(path):

    prs = Presentation(path)

    extracted = []

    for slide_num, slide in enumerate(prs.slides):

        slide_text = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                slide_text.append(shape.text)

        text = clean_text("\n".join(slide_text))

        if len(text.strip()) < 30:
            continue

        extracted.append({
            "text": text,
            "source": os.path.basename(path),
            "page": slide_num + 1
        })

    return extracted